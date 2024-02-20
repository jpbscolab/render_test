from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

from pathlib import Path
import json

current = Path()
template_file_path = current / "files" / "jpowergroup.pptx"

def createSamplePptx():

    # Create a presentation object
    prs = Presentation(template_file_path)

    # Grand title
    presentation_title = "生成AIについて"

    # Define slide titles
    slide_titles = [
        "生成AIのLLMが得意なこと、不得意なこと",
        "生成AIのLLMはパターンをみつけたりパターンにあてはめるのが得意",
        "プロンプトの組み立て方"
    ]

    # Define slide content
    slide_contents = [
        "得意なこと：テキスト生成、言語理解、翻訳\n不得意なこと：感情を理解する、画像や音声を解析する",
        "LLMは大量のデータからパターンを学習し、新しいシナリオにこれらのパターンを適用することが得意です。",
        "良いプロンプトは明確で具体的です。目的に応じて、必要な情報や質問を組み込みましょう。"
    ]

    # Define slide notes (to be added in the speaker notes section)
    slide_notes = [
        "このスライドでは、生成AIのLLMが得意とする分野と苦手とする分野について説明します。",
        "LLMがパターンを見つけ、それを応用する能力について詳しく述べます。",
        "プロンプトの組み立て方について、有効な例とヒントを提供します。"
    ]

    # Define keywords to emphasize
    keywords = [
        ["テキスト生成", "言語理解", "翻訳", "感情", "画像", "音声"],
        ["パターン", "学習", "適用"],
        ["プロンプト", "明確", "具体的", "情報", "質問"]
    ]

    json_str = json.dumps({'presentation_title':presentation_title, 'slide_titles':slide_titles, 'slide_contents':slide_contents, 'slide_notes':slide_notes, 'keywords':keywords})

    #file_name = createPptx(slide_titles, slide_contents, slide_notes, keywords)
    file_name = createPptxFromJson(json_str)
    return file_name

def addRunsToTextFrame(textFrame, text, keywords):
    # keywords に該当するものがあるかを
    # textの中から検索していき、該当する箇所の直前と該当部分を切り出しリストに吐き出す
    # 残りの部分にも同じ処理をする
    runs = []
    resttext = text
    paragraph = textFrame.paragraphs[0]
    while True:
        found = len(resttext)
        foundKeyword = ""
        for keyword in (keywords):
            tmp = resttext.find(keyword)
            if tmp >= 0 and tmp < found:
                found = tmp
                foundKeyword = keyword
        if found == len(resttext):
            runs += [resttext]
            for _run in runs:
                run = paragraph.add_run()
                run.text = _run
            return runs
        runs += [resttext[0:found], foundKeyword]
        resttext = resttext[found + len(foundKeyword):]

def createPptx(presentation_title, slide_titles, slide_contents, slide_notes, keywords):

    # Create a presentation object
    prs = Presentation(template_file_path)

    # Define colors
    # bg_color = RGBColor(0, 32, 96)  # Dark blue background for good visibility
    title_color = RGBColor(64, 128, 255)
    font_color = RGBColor(0, 0, 0)  # White font for contrast
    accent_color = RGBColor(255, 96, 0)  # Accent color for key elements

    # Define font size
    font_size = Pt(24)

    # Presentation title
    prs.slides[0].shapes[0].text_frame.text = presentation_title

    # Helper function to add a slide
    def add_slide(title, content, note, keywords):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        title_shape = slide.shapes.title
        title_shape.text = title

        # Set the title style
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.color.rgb = title_color

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(9.2), Inches(5))
        content_frame = content_box.text_frame
        # content_frame.text = textToRunTexts(content, keywords)
        addRunsToTextFrame(content_frame, content, keywords)
        content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        content_frame.word_wrap = True

        # Style the content
        for i, paragraph in enumerate(content_frame.paragraphs):
            for run in paragraph.runs:
                # print(run.text)
                run.font.size = font_size
                run.font.color.rgb = font_color
                # If the content line is a keyword, change the font color to accent color
                if any(keyword in run.text for keyword in keywords):
                    run.font.color.rgb = accent_color

        # Add note to the slide
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

        # Set the background color of the slide
        background = slide.background
        fill = background.fill
        fill.solid()
        # fill.fore_color.rgb = bg_color

    # Add slides to the presentation
    for title, content, note, kw in zip(slide_titles, slide_contents, slide_notes, keywords):
        add_slide(title, content, note, kw)

    # Save the presentation to a file
    file_name = 'presentation.pptx'
    file_path = f'./files/{file_name}'
    prs.save(file_path)
    return file_name

def createPptxFromJson(json_str):
    json_dict = json.loads(json_str)
    presentation_title = ''
    slide_titles = []
    slide_contents = []
    slide_notes = []
    keywords = []
    if 'presentation_title' in json_dict:
        presentation_title = json_dict['presentation_title']
    if 'slide_titles' in json_dict:
        slide_titles = json_dict['slide_titles']
    if 'slide_contents' in json_dict:
        slide_contents = json_dict['slide_contents']
    if 'slide_notes' in json_dict:
        slide_notes = json_dict['slide_notes']
    if 'keywords' in json_dict:
        keywords = json_dict['keywords']

    file_name = createPptx(presentation_title, slide_titles, slide_contents, slide_notes, keywords)
    return file_name
